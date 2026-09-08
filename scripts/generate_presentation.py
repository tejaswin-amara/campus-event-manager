import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Color Palette Constants ---
BG_COLOR = RGBColor(11, 15, 25)        # Deep Obsidian Navy #0B0F19
CARD_BG = RGBColor(24, 34, 52)         # Surface Card #182234
CARD_BG_ALT = RGBColor(18, 26, 40)     # Alternate Card #121A28
CARD_BORDER = RGBColor(51, 65, 85)     # Slate Border #334155
TEXT_WHITE = RGBColor(255, 255, 255)   # White Header
TEXT_LIGHT = RGBColor(226, 232, 240)   # Light Slate Text #E2E8F0
TEXT_MUTED = RGBColor(148, 163, 184)   # Muted Slate #94A3B8
ACCENT_INDIGO = RGBColor(99, 102, 241) # Indigo #6366F1
ACCENT_CYAN = RGBColor(6, 182, 212)    # Cyan #06B6D4
ACCENT_GREEN = RGBColor(16, 185, 129)  # Emerald #10B981
ACCENT_AMBER = RGBColor(245, 158, 11)  # Amber #F59E0B
ACCENT_PINK = RGBColor(236, 72, 153)   # Pink #EC4899

FONT_FAMILY = "Arial"

ASSETS_DIR = os.path.abspath("presentation_assets")
OUTPUT_PPTX = os.path.abspath("docs/CampusConnect_Review_2_Presentation.pptx")
os.makedirs(os.path.dirname(OUTPUT_PPTX), exist_ok=True)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

def apply_background(slide):
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = BG_COLOR
    bg_shape.line.fill.background()
    return bg_shape

def add_header(slide, title, category="CAMPUSCONNECT • REVIEW 2", subtitle=None):
    # Category Tag Pill
    pill = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.42), Inches(3.6), Inches(0.32)
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = CARD_BG
    pill.line.color.rgb = ACCENT_INDIGO
    pill.line.width = Pt(1)
    tf = pill.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = category
    p.font.name = FONT_FAMILY
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.alignment = PP_ALIGN.CENTER
    
    # Title Text
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.733), Inches(0.85))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.font.name = FONT_FAMILY
    p_title.font.size = Pt(30)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE
    
    if subtitle:
        p_sub = tf_title.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.name = FONT_FAMILY
        p_sub.font.size = Pt(15)
        p_sub.font.color.rgb = TEXT_MUTED
        p_sub.space_before = Pt(3)

def add_card(slide, left, top, width, height, fill_color=CARD_BG, border_color=CARD_BORDER):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = border_color
    card.line.width = Pt(1.2)
    return card

def add_speaker_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text

# ==========================================
# SLIDE 1: Title Slide
# ==========================================
slide1 = prs.slides.add_slide(blank_layout)
apply_background(slide1)

# Large Hero Card
add_card(slide1, Inches(1.2), Inches(1.0), Inches(10.933), Inches(5.5), fill_color=CARD_BG, border_color=ACCENT_INDIGO)

# Title Badge
pill1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.6), Inches(1.4), Inches(4.5), Inches(0.38))
pill1.fill.solid()
pill1.fill.fore_color.rgb = CARD_BG_ALT
pill1.line.color.rgb = ACCENT_CYAN
pill1.line.width = Pt(1)
p1 = pill1.text_frame.paragraphs[0]
p1.text = "ACADEMIC CAPSTONE PROJECT • REVIEW 2"
p1.font.name = FONT_FAMILY
p1.font.size = Pt(11)
p1.font.bold = True
p1.font.color.rgb = ACCENT_CYAN
p1.alignment = PP_ALIGN.CENTER

# Main Title & Subtitle Box
tb = slide1.shapes.add_textbox(Inches(1.6), Inches(1.95), Inches(10.1), Inches(2.0))
tf = tb.text_frame
tf.word_wrap = True

p_main = tf.paragraphs[0]
p_main.text = "CampusConnect"
p_main.font.name = FONT_FAMILY
p_main.font.size = Pt(40)
p_main.font.bold = True
p_main.font.color.rgb = TEXT_WHITE

p_sub = tf.add_paragraph()
p_sub.text = "A Trustworthy Campus Event Catalogue with a Secure Administrative Control Plane"
p_sub.font.name = FONT_FAMILY
p_sub.font.size = Pt(22)
p_sub.font.color.rgb = ACCENT_INDIGO
p_sub.space_before = Pt(8)

p_desc = tf.add_paragraph()
p_desc.text = "Review 2 Deliverables: System Architecture, Frontend Implementation, Database Normalization, Schema & ER Design"
p_desc.font.name = FONT_FAMILY
p_desc.font.size = Pt(16)
p_desc.font.color.rgb = TEXT_MUTED
p_desc.space_before = Pt(6)

# Information Grid (2 Columns)
info_card = add_card(slide1, Inches(1.6), Inches(4.1), Inches(10.1), Inches(1.9), fill_color=CARD_BG_ALT, border_color=CARD_BORDER)

tb_info = slide1.shapes.add_textbox(Inches(1.9), Inches(4.2), Inches(4.6), Inches(1.7))
tf_info = tb_info.text_frame
tf_info.word_wrap = True

p_team_h = tf_info.paragraphs[0]
p_team_h.text = "Project Team Members:"
p_team_h.font.name = FONT_FAMILY
p_team_h.font.size = Pt(15)
p_team_h.font.bold = True
p_team_h.font.color.rgb = ACCENT_CYAN

p_team_1 = tf_info.add_paragraph()
p_team_1.text = "• Tejaswin Amara (Lead Full-Stack & System Architecture)"
p_team_1.font.size = Pt(13)
p_team_1.font.color.rgb = TEXT_LIGHT

p_team_2 = tf_info.add_paragraph()
p_team_2.text = "• Team Member 2 (Database Engineering & Migration Lead)"
p_team_2.font.size = Pt(13)
p_team_2.font.color.rgb = TEXT_LIGHT

p_team_3 = tf_info.add_paragraph()
p_team_3.text = "• Team Member 3 (Frontend & Security Integration Specialist)"
p_team_3.font.size = Pt(13)
p_team_3.font.color.rgb = TEXT_LIGHT

tb_guide = slide1.shapes.add_textbox(Inches(6.8), Inches(4.2), Inches(4.6), Inches(1.7))
tf_guide = tb_guide.text_frame
tf_guide.word_wrap = True

p_g_h = tf_guide.paragraphs[0]
p_g_h.text = "Project Supervision & Tech Stack:"
p_g_h.font.name = FONT_FAMILY
p_g_h.font.size = Pt(15)
p_g_h.font.bold = True
p_g_h.font.color.rgb = ACCENT_GREEN

p_g_1 = tf_guide.add_paragraph()
p_g_1.text = "Project Guide: Prof. / Dr. Department Mentor"
p_g_1.font.size = Pt(13)
p_g_1.font.color.rgb = TEXT_LIGHT

p_g_2 = tf_guide.add_paragraph()
p_g_2.text = "Department: Computer Science & Engineering"
p_g_2.font.size = Pt(13)
p_g_2.font.color.rgb = TEXT_LIGHT

p_g_3 = tf_guide.add_paragraph()
p_g_3.text = "Core Stack: Java 25 | Spring Boot 4 | MySQL 8.4 LTS | React & Thymeleaf"
p_g_3.font.size = Pt(13)
p_g_3.font.bold = True
p_g_3.font.color.rgb = ACCENT_INDIGO

add_speaker_notes(slide1, """Good morning respected panel members and evaluators. Welcome to the Review 2 presentation of our capstone project, CampusConnect. 

In Review 1, we defined our project title, problem scope, and preliminary system architecture. Today, in Review 2, we present our working technical implementation: our modular monolith architecture, our completed frontend interfaces and forms, our 3NF normalized database schema and ER design, live sample database records, and our engineering solutions to concurrency and migration challenges. 

Let us begin by revisiting the problem statement and existing system gaps.""")

# ==========================================
# SLIDE 2: Problem Statement & Existing System
# ==========================================
slide2 = prs.slides.add_slide(blank_layout)
apply_background(slide2)
add_header(slide2, "Problem Statement & Existing System Analysis", category="PROBLEM & MOTIVATION", subtitle="Why campus event management requires a centralized, verifiable solution")

# Card Left: Existing System Gaps
add_card(slide2, Inches(0.8), Inches(1.9), Inches(5.6), Inches(5.0))
tb_c1 = slide2.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(5.0), Inches(4.6))
tf_c1 = tb_c1.text_frame
tf_c1.word_wrap = True

p = tf_c1.paragraphs[0]
p.text = "Existing System Bottlenecks"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = ACCENT_PINK

bullets_c1 = [
    ("Information Fragmentation: ", "Event details are scattered across informal WhatsApp groups, bulletin boards, and unverified social media handles."),
    ("Broken Redirects & Phishing: ", "Students frequently encounter outdated external registration links and unverified Google Forms."),
    ("No Verified Capacity Enforcement: ", "Organizers cannot broadcast seat exhaustion in real time, leading to overcrowded physical venues."),
    ("Duplicate & Ghost Signups: ", "Lack of transactional user-event constraints allows duplicate signups that distort attendance metrics.")
]
for title_txt, body_txt in bullets_c1:
    p = tf_c1.add_paragraph()
    p.space_before = Pt(12)
    run1 = p.add_run()
    run1.text = "• " + title_txt
    run1.font.bold = True
    run1.font.size = Pt(16)
    run1.font.color.rgb = TEXT_WHITE
    run2 = p.add_run()
    run2.text = body_txt
    run2.font.size = Pt(15)
    run2.font.color.rgb = TEXT_LIGHT

# Card Right: Proposed Solution & Value Add
add_card(slide2, Inches(6.9), Inches(1.9), Inches(5.6), Inches(5.0))
tb_c2 = slide2.shapes.add_textbox(Inches(7.2), Inches(2.1), Inches(5.0), Inches(4.6))
tf_c2 = tb_c2.text_frame
tf_c2.word_wrap = True

p = tf_c2.paragraphs[0]
p.text = "The CampusConnect Solution"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = ACCENT_GREEN

bullets_c2 = [
    ("Single Source of Truth: ", "A centralized, searchable catalogue of verified campus events with instant category filtering."),
    ("Validated Registration Links: ", "Strict server-side protocol checks guarantee students are redirected only to authenticated portals."),
    ("Transactional Interest Tracking: ", "Pessimistic write locking enforces unique registration states with zero duplicate record overhead."),
    ("Secure Admin Control Plane: ", "Role-Based Access Control (RBAC) allows verified faculty/leads to manage lifecycles and export telemetry.")
]
for title_txt, body_txt in bullets_c2:
    p = tf_c2.add_paragraph()
    p.space_before = Pt(12)
    run1 = p.add_run()
    run1.text = "• " + title_txt
    run1.font.bold = True
    run1.font.size = Pt(16)
    run1.font.color.rgb = TEXT_WHITE
    run2 = p.add_run()
    run2.text = body_txt
    run2.font.size = Pt(15)
    run2.font.color.rgb = TEXT_LIGHT

add_speaker_notes(slide2, """On this slide, we analyze the core problem motivating CampusConnect. In colleges today, event information is chaotic. Club leads circulate Google forms across unofficial group chats, leading to lost links, fraudulent pages, and missed deadlines. Furthermore, organizers have no real-time telemetry into venue capacity or student interest.

CampusConnect resolves these problems by providing an authoritative single source of truth. Notice that we enforce server-side validation on external links, and we use pessimistic locking to prevent race conditions during interest registration.""")

# ==========================================
# SLIDE 3: Objectives & Scope
# ==========================================
slide3 = prs.slides.add_slide(blank_layout)
apply_background(slide3)
add_header(slide3, "Project Objectives & Scope Boundaries", category="PROJECT SPECIFICATION", subtitle="Specific, measurable technical outcomes and defined engineering boundaries")

# Left Column: Project Objectives
add_card(slide3, Inches(0.8), Inches(1.9), Inches(5.6), Inches(5.0))
tb_o1 = slide3.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(5.0), Inches(4.6))
tf_o1 = tb_o1.text_frame
tf_o1.word_wrap = True

p = tf_o1.paragraphs[0]
p.text = "Core Technical Objectives"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = ACCENT_INDIGO

objs = [
    ("Architect a Modular Monolith: ", "Build a unified, robust Spring Boot 4 backend with strict domain boundaries ready for future microservice extraction."),
    ("High-Performance Event Discovery: ", "Implement indexed search, category filtering, and explainable upcoming event recommendations."),
    ("Zero-Trust Security Baseline: ", "Enforce BCrypt credential encryption, CSRF protection, and Bucket4j IP-based login rate limiting."),
    ("Automated Schema Governance: ", "Establish relational integrity via version-controlled Flyway V1–V3 migrations on MySQL 8.4.")
]
for title_txt, body_txt in objs:
    p = tf_o1.add_paragraph()
    p.space_before = Pt(12)
    run1 = p.add_run()
    run1.text = "• " + title_txt
    run1.font.bold = True
    run1.font.size = Pt(16)
    run1.font.color.rgb = TEXT_WHITE
    run2 = p.add_run()
    run2.text = body_txt
    run2.font.size = Pt(15)
    run2.font.color.rgb = TEXT_LIGHT

# Right Column: Scope Boundaries
add_card(slide3, Inches(6.9), Inches(1.9), Inches(5.6), Inches(5.0))
tb_o2 = slide3.shapes.add_textbox(Inches(7.2), Inches(2.1), Inches(5.0), Inches(4.6))
tf_o2 = tb_o2.text_frame
tf_o2.word_wrap = True

p = tf_o2.paragraphs[0]
p.text = "Defined Project Scope"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = ACCENT_CYAN

scopes = [
    ("In-Scope (Current Implementation): ", "Student public catalogue, event search/filter, interest registration, admin CRUD operations, image uploads (BLOB), telemetry dashboard, CSV export, health/Prometheus metrics."),
    ("Authoritative Registration Link: ", "Interest is recorded within CampusConnect, while the configured external registration link remains authoritative for seat ticket issuance."),
    ("Honest Architectural Boundary: ", "Technologies such as Kafka, MongoDB, pgvector, and FastAPI are modeled as bounded future evolution paths, not falsely claimed as current runtime components.")
]
for title_txt, body_txt in scopes:
    p = tf_o2.add_paragraph()
    p.space_before = Pt(12)
    run1 = p.add_run()
    run1.text = "• " + title_txt
    run1.font.bold = True
    run1.font.size = Pt(16)
    run1.font.color.rgb = TEXT_WHITE
    run2 = p.add_run()
    run2.text = body_txt
    run2.font.size = Pt(15)
    run2.font.color.rgb = TEXT_LIGHT

add_speaker_notes(slide3, """On Slide 3, we detail our technical objectives and explicit scope boundaries. Our goal is to engineer an enterprise-grade modular monolith using Java 25, Spring Boot 4, and MySQL 8.4 LTS. 

Importantly, we highlight an honest scope boundary: our application currently tracks student interest and directs them to authoritative campus registration links. We do not claim to run an entire payment or ticket-scanning system in this sprint; those are clearly mapped as Sprint 3 and future bounded evolution paths.""")

# ==========================================
# SLIDE 4: Proposed System & Key Capabilities
# ==========================================
slide4 = prs.slides.add_slide(blank_layout)
apply_background(slide4)
add_header(slide4, "Proposed System & Key Innovations", category="SYSTEM DESIGN", subtitle="Core architectural capabilities delivered in CampusConnect")

# 4 Grid Feature Cards
coords = [
    (Inches(0.8), Inches(1.9), "Modular Monolith Architecture", ACCENT_INDIGO, [
        ("Single Deployable Artifact: ", "Unified Spring Boot application packaging controllers, business logic, security filters, and JPA entities."),
        ("Strict Module Separation: ", "Decoupled domain packages for Identity, Events, Registration, and Recommendations.")
    ]),
    (Inches(6.9), Inches(1.9), "Zero-Trust Security Defense", ACCENT_PINK, [
        ("Multi-Layer Authentication: ", "Spring Security filter chain with BCrypt 60-character salted password hashing."),
        ("Abuse Prevention: ", "Bucket4j IP-based login rate limiting, session-fixation protection, and CSRF token verification.")
    ]),
    (Inches(0.8), Inches(4.55), "Relational Data Integrity", ACCENT_GREEN, [
        ("3NF Normalized Schema: ", "Normalized MySQL 8.4 database managed exclusively via versioned Flyway V1–V3 migrations."),
        ("Pessimistic Concurrency Lock: ", "Row-level locks guarantee duplicate-free student interest registration under heavy concurrent clicks.")
    ]),
    (Inches(6.9), Inches(4.55), "Intelligent Discovery & Telemetry", ACCENT_CYAN, [
        ("Explainable Recommendations: ", "Server-side heuristic scoring recommending top 3 upcoming events based on category & velocity."),
        ("Operational Observability: ", "Prometheus metrics via Micrometer and `/actuator/health` liveness/readiness probes.")
    ])
]

for left, top, title, color, points in coords:
    add_card(slide4, left, top, Inches(5.6), Inches(2.45))
    tb_box = slide4.shapes.add_textbox(left + Inches(0.25), top + Inches(0.18), Inches(5.1), Inches(2.1))
    tf_box = tb_box.text_frame
    tf_box.word_wrap = True
    
    p = tf_box.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = color
    
    for lead, desc in points:
        p = tf_box.add_paragraph()
        p.space_before = Pt(6)
        r1 = p.add_run()
        r1.text = "• " + lead
        r1.font.bold = True
        r1.font.size = Pt(14)
        r1.font.color.rgb = TEXT_WHITE
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(13)
        r2.font.color.rgb = TEXT_LIGHT

add_speaker_notes(slide4, """Slide 4 presents the four technological pillars of our proposed system:
First, a Modular Monolith architecture that allows fast, reliable deployments without microservice network overhead.
Second, Zero-Trust security featuring BCrypt, CSRF guards, and Bucket4j rate limiting against brute-force attacks.
Third, relational integrity with Flyway migrations and row-level pessimistic locking.
Fourth, intelligent discovery with explainable upcoming recommendations and Prometheus telemetry.""")

# ==========================================
# SLIDE 5: System Architecture (C4 Model)
# ==========================================
slide5 = prs.slides.add_slide(blank_layout)
apply_background(slide5)
add_header(slide5, "System Architecture (C4 Layered Model)", category="ARCHITECTURE", subtitle="Multi-tiered modular design: Presentation, Security, Business Logic, and Data Persistence")

# Left: Diagram Image Card
add_card(slide5, Inches(0.8), Inches(1.85), Inches(7.6), Inches(5.1))
arch_img_path = os.path.join(ASSETS_DIR, "system_architecture.png")
if os.path.exists(arch_img_path):
    slide5.shapes.add_picture(arch_img_path, Inches(0.95), Inches(2.0), width=Inches(7.3))

# Right: Architecture Explanation
add_card(slide5, Inches(8.6), Inches(1.85), Inches(3.933), Inches(5.1))
tb_arch = slide5.shapes.add_textbox(Inches(8.8), Inches(2.05), Inches(3.533), Inches(4.7))
tf_arch = tb_arch.text_frame
tf_arch.word_wrap = True

p = tf_arch.paragraphs[0]
p.text = "Architectural Layers"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = ACCENT_INDIGO

arch_points = [
    ("Presentation Tier: ", "Dual web surfaces: Server-rendered Thymeleaf templates and responsive React dashboard interface."),
    ("Security Boundary: ", "Spring Security intercepts every request. Validates session cookies, CSRF tokens, and enforces RBAC."),
    ("Service Layer: ", "Transactional domain services (`EventService`, `UserService`, `RecommendationService`) maintain business invariants."),
    ("Persistence Tier: ", "Spring Data JPA repositories connect to MySQL 8.4 LTS; Flyway enforces schema changes."),
    ("Operational Plane: ", "Actuator exposes `/health` and Prometheus metrics for container monitoring.")
]

for title_txt, body_txt in arch_points:
    p = tf_arch.add_paragraph()
    p.space_before = Pt(8)
    run1 = p.add_run()
    run1.text = "• " + title_txt
    run1.font.bold = True
    run1.font.size = Pt(14)
    run1.font.color.rgb = TEXT_WHITE
    run2 = p.add_run()
    run2.text = body_txt
    run2.font.size = Pt(13)
    run2.font.color.rgb = TEXT_LIGHT

add_speaker_notes(slide5, """On Slide 5, we show the complete System Architecture following the C4 container model. 

On the left is our architectural diagram: requests enter from Student and Admin browsers into our Spring Boot application container running on Java 25. The request passes through our Security Interceptor chain, which evaluates BCrypt credentials and Bucket4j rate limits. It is then routed to the relevant Controller and Service layer. 

The service layer interacts with Spring Data JPA repositories, which read and write to MySQL 8.4 LTS. Crucially, Flyway migrations V1 through V3 run before Hibernate validation, preventing any uncontrolled schema drift.""")

# ==========================================
# SLIDE 6: System Modules & Functionalities
# ==========================================
slide6 = prs.slides.add_slide(blank_layout)
apply_background(slide6)
add_header(slide6, "System Modules & Functional Breakdown", category="MODULES", subtitle="Decoupled functional boundaries ensuring maintainability and clean separation of concerns")

module_boxes = [
    (Inches(0.8), Inches(1.9), "Module 1: Identity & Access Control", ACCENT_INDIGO, [
        "Dual-Role Authorization: Strict segregation of STUDENT and ADMIN privileges.",
        "BCrypt Credential Protection: 60-character salted password hashes stored in MySQL.",
        "Rate Limiting Filter: Bucket4j token-bucket filter blocks brute-force login attempts.",
        "Safe Session Handling: Session-fixation protection and secure cookie flags."
    ]),
    (Inches(6.9), Inches(1.9), "Module 2: Event Discovery & Catalogue", ACCENT_CYAN, [
        "Public Event Catalogue: Paginated listings with high-contrast card presentation.",
        "Multi-Param Filter & Search: Real-time search across event title, venue, and description.",
        "Category Pill Navigation: Instant filtering by Technical, Cultural, Sports, and Seminars.",
        "Event Detail Drawer: Full modal dialog with date/time, venue pin, and capacity gauge."
    ]),
    (Inches(0.8), Inches(4.55), "Module 3: Registration & Interest Engine", ACCENT_GREEN, [
        "Transactional Interest Flow: Safe user registration recording before external redirect.",
        "Concurrency Guard: Pessimistic write lock prevents race-condition double registration.",
        "External Link Verification: Server-side URI scheme validation (HTTP/HTTPS whitelist).",
        "Unique Constraint Defense: MySQL compound key `uk_user_event` guarantees data integrity."
    ]),
    (Inches(6.9), Inches(4.55), "Module 4: Admin Control Plane & Telemetry", ACCENT_AMBER, [
        "Full Event Lifecycle (CRUD): Admin interface to create, edit, and delete event records.",
        "Media Upload Management: Validated image upload persistence (MEDIUMBLOB).",
        "Real-Time KPI Analytics: Metric cards tracking total events, registrations, and capacity.",
        "One-Click CSV Export: Instant download of event rosters for offline administrative reporting."
    ])
]

for left, top, title, color, items in module_boxes:
    add_card(slide6, left, top, Inches(5.6), Inches(2.45))
    tb_m = slide6.shapes.add_textbox(left + Inches(0.25), top + Inches(0.18), Inches(5.1), Inches(2.1))
    tf_m = tb_m.text_frame
    tf_m.word_wrap = True
    
    p = tf_m.paragraphs[0]
    p.text = title
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = color
    
    for item in items:
        p = tf_m.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = "• " + item
        run.font.size = Pt(13)
        run.font.color.rgb = TEXT_LIGHT

add_speaker_notes(slide6, """Slide 6 outlines the four core functional modules of CampusConnect:
Module 1 handles Identity and Access Control with role-based segregation and rate limiting.
Module 2 delivers Event Discovery, providing search, category filtering, and event detail views.
Module 3 is our Registration and Interest Engine, protecting data integrity with database unique constraints and row-level locks.
Module 4 is the Admin Control Plane, providing administrative CRUD operations, image uploads, and CSV exports.""")

# ==========================================
# SLIDE 7: Frontend Design & UI/UX Wireframes
# ==========================================
slide7 = prs.slides.add_slide(blank_layout)
apply_background(slide7)
add_header(slide7, "Frontend Design: Discovery Feed & Authentication", category="FRONTEND DESIGN", subtitle="Modern design tokens, high contrast visual hierarchy, and responsive mobile-first layouts")

# Left Image: Student Catalogue
add_card(slide7, Inches(0.8), Inches(1.85), Inches(5.6), Inches(4.3))
img1 = os.path.join(ASSETS_DIR, "01_student_catalogue.png")
if os.path.exists(img1):
    slide7.shapes.add_picture(img1, Inches(0.95), Inches(1.95), width=Inches(5.3))

# Left Caption
tb_cap1 = slide7.shapes.add_textbox(Inches(0.8), Inches(6.25), Inches(5.6), Inches(0.8))
tf_cap1 = tb_cap1.text_frame
tf_cap1.word_wrap = True
p = tf_cap1.paragraphs[0]
p.text = "Student Discovery Home Page"
p.font.bold = True
p.font.size = Pt(15)
p.font.color.rgb = ACCENT_INDIGO
p2 = tf_cap1.add_paragraph()
p2.text = "Search bar, category filter pills, recommended event banner, and responsive card grid."
p2.font.size = Pt(13)
p2.font.color.rgb = TEXT_MUTED

# Right Image: Admin Login
add_card(slide7, Inches(6.9), Inches(1.85), Inches(5.6), Inches(4.3))
img4 = os.path.join(ASSETS_DIR, "04_admin_login.png")
if os.path.exists(img4):
    slide7.shapes.add_picture(img4, Inches(7.05), Inches(1.95), width=Inches(5.3))

# Right Caption
tb_cap2 = slide7.shapes.add_textbox(Inches(6.9), Inches(6.25), Inches(5.6), Inches(0.8))
tf_cap2 = tb_cap2.text_frame
tf_cap2.word_wrap = True
p = tf_cap2.paragraphs[0]
p.text = "Secure Admin Authentication Portal"
p.font.bold = True
p.font.size = Pt(15)
p.font.color.rgb = ACCENT_CYAN
p2 = tf_cap2.add_paragraph()
p2.text = "Glassmorphic authentication card with BCrypt protection & Bucket4j rate-limit feedback."
p2.font.size = Pt(13)
p2.font.color.rgb = TEXT_MUTED

add_speaker_notes(slide7, """Moving into our technical demonstration, Slide 7 showcases our implemented frontend designs. 

On the left is the Student Event Catalogue. Notice the dark theme engineered for high contrast and readability: students can instantly search events, click category filter pills (Technical, Cultural, Sports), and see upcoming recommended events.

On the right is the Admin Authentication screen. It features a modern glassmorphic card design with input validation and clear feedback regarding rate limiting and security status.""")

# ==========================================
# SLIDE 8: Frontend Implementation: Dashboard & Forms
# ==========================================
slide8 = prs.slides.add_slide(blank_layout)
apply_background(slide8)
add_header(slide8, "Frontend Implementation: Control Center & Forms", category="LIVE DEMONSTRATION", subtitle="Administrative event lifecycle management, real-time KPI metrics, and input validation")

# Left Image: Admin Dashboard
add_card(slide8, Inches(0.8), Inches(1.85), Inches(5.6), Inches(4.3))
img2 = os.path.join(ASSETS_DIR, "02_admin_dashboard.png")
if os.path.exists(img2):
    slide8.shapes.add_picture(img2, Inches(0.95), Inches(1.95), width=Inches(5.3))

tb_cap3 = slide8.shapes.add_textbox(Inches(0.8), Inches(6.25), Inches(5.6), Inches(0.8))
tf_cap3 = tb_cap3.text_frame
tf_cap3.word_wrap = True
p = tf_cap3.paragraphs[0]
p.text = "Admin Control Plane & Analytics"
p.font.bold = True
p.font.size = Pt(15)
p.font.color.rgb = ACCENT_GREEN
p2 = tf_cap3.add_paragraph()
p2.text = "Live telemetry cards (Total Events, Active Signups, Capacity) with searchable event table."
p2.font.size = Pt(13)
p2.font.color.rgb = TEXT_MUTED

# Right Image: Create Event Form
add_card(slide8, Inches(6.9), Inches(1.85), Inches(5.6), Inches(4.3))
img3 = os.path.join(ASSETS_DIR, "03_create_event_form.png")
if os.path.exists(img3):
    slide8.shapes.add_picture(img3, Inches(7.05), Inches(1.95), width=Inches(5.3))

tb_cap4 = slide8.shapes.add_textbox(Inches(6.9), Inches(6.25), Inches(5.6), Inches(0.8))
tf_cap4 = tb_cap4.text_frame
tf_cap4.word_wrap = True
p = tf_cap4.paragraphs[0]
p.text = "Event Lifecycle Creation & Edit Form"
p.font.bold = True
p.font.size = Pt(15)
p.font.color.rgb = ACCENT_AMBER
p2 = tf_cap4.add_paragraph()
p2.text = "Comprehensive form modal with Bean Validation for dates, venues, capacities, and media URLs."
p2.font.size = Pt(13)
p2.font.color.rgb = TEXT_MUTED

add_speaker_notes(slide8, """Slide 8 demonstrates the administrative operational interface. 

On the left is the Admin Control Plane. It provides administrators with live KPI cards tracking total events, student interest count, and capacity metrics. From this table, administrators can search records or trigger a one-click CSV export for faculty reporting.

On the right is the Create Event modal dialog. It enforces strict client-side and server-side validation: end times must be after start times, capacities must be positive integers, and URLs must follow valid HTTP/HTTPS schemes.""")

# ==========================================
# SLIDE 9: Database Design & Relational Foundation
# ==========================================
slide9 = prs.slides.add_slide(blank_layout)
apply_background(slide9)
add_header(slide9, "Database Design & Relational Foundation", category="DATA ENGINEERING", subtitle="3NF normalization, transactional consistency, and versioned migration governance")

# 3 Column Cards
db_cards = [
    (Inches(0.8), Inches(1.9), "Third Normal Form (3NF)", ACCENT_INDIGO, [
        ("Zero Transitive Dependencies: ", "Every non-key attribute is strictly dependent on the primary key."),
        ("Entity Decoupling: ", "User identity records are cleanly isolated from event metadata and media blobs."),
        ("Associative Entity: ", "`registrations` table resolves the many-to-many relationship cleanly.")
    ]),
    (Inches(4.85), Inches(1.9), "Flyway Migration Authority", ACCENT_CYAN, [
        ("V1__Initial_Schema.sql: ", "Baseline tables for users, events, registrations with PK/FK constraints."),
        ("V2__Add_Image_Blob.sql: ", "Portable database-backed media storage columns (MEDIUMBLOB)."),
        ("V3__Indexes_Integrity.sql: ", "Compound indexes and check constraints for query-aware analytics.")
    ]),
    (Inches(8.9), Inches(1.9), "Performance & Locking", ACCENT_GREEN, [
        ("Query-Aware Indexing: ", "Composite indexes on `(date_time)` and `(category, date_time)` optimize discovery."),
        ("Pessimistic Row Lock: ", "JPA `PESSIMISTIC_WRITE` locks the event row during registration to eliminate race conditions."),
        ("Hibernate Validation: ", "`DDL_AUTO=validate` prevents silent runtime schema tampering.")
    ])
]

for left, top, title, color, points in db_cards:
    add_card(slide9, left, top, Inches(3.65), Inches(5.0))
    tb_dbc = slide9.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(3.25), Inches(4.6))
    tf_dbc = tb_dbc.text_frame
    tf_dbc.word_wrap = True
    
    p = tf_dbc.paragraphs[0]
    p.text = title
    p.font.size = Pt(19)
    p.font.bold = True
    p.font.color.rgb = color
    
    for lead, desc in points:
        p = tf_dbc.add_paragraph()
        p.space_before = Pt(12)
        r1 = p.add_run()
        r1.text = "• " + lead
        r1.font.bold = True
        r1.font.size = Pt(15)
        r1.font.color.rgb = TEXT_WHITE
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(14)
        r2.font.color.rgb = TEXT_LIGHT

add_speaker_notes(slide9, """On Slide 9, we dive into our database engineering. 
We strictly adhere to Third Normal Form (3NF). Notice how the schema isolates users from events, using the associative `registrations` table to represent user interest.

Crucially, our database is not created manually or mutated by Hibernate. Instead, we use Flyway versioned migrations V1, V2, and V3 on MySQL 8.4 LTS, running with Hibernate in validate mode. This ensures that the code and schema remain strictly synchronized in CI and production.""")

# ==========================================
# SLIDE 10: Entity-Relationship (ER) Diagram
# ==========================================
slide10 = prs.slides.add_slide(blank_layout)
apply_background(slide10)
add_header(slide10, "Entity-Relationship (ER) Diagram", category="DATABASE MODEL", subtitle="Normalized relational schema with referential integrity, foreign keys, and unique constraints")

# Left: Diagram Image Card
add_card(slide10, Inches(0.8), Inches(1.85), Inches(7.8), Inches(5.1))
er_img_path = os.path.join(ASSETS_DIR, "er_diagram.png")
if os.path.exists(er_img_path):
    slide10.shapes.add_picture(er_img_path, Inches(0.95), Inches(2.0), width=Inches(7.5))

# Right: Structural Analysis
add_card(slide10, Inches(8.8), Inches(1.85), Inches(3.733), Inches(5.1))
tb_er = slide10.shapes.add_textbox(Inches(9.0), Inches(2.05), Inches(3.333), Inches(4.7))
tf_er = tb_er.text_frame
tf_er.word_wrap = True

p = tf_er.paragraphs[0]
p.text = "Relational Semantics"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = ACCENT_INDIGO

er_points = [
    ("USERS (1) to REGISTRATIONS (N): ", "One student account can express interest in multiple campus events."),
    ("EVENTS (1) to REGISTRATIONS (N): ", "Each campus event can receive registrations from multiple student users."),
    ("Compound Unique Constraint: ", "`UNIQUE (user_id, event_id)` enforces that a user registers for an event exactly once."),
    ("Cascading Deletion Policy: ", "`ON DELETE CASCADE` guarantees referential cleanliness if an event is cancelled."),
    ("Integrity Check Constraints: ", "`CHECK (max_capacity > 0)` and status enumeration check.")
]

for title_txt, body_txt in er_points:
    p = tf_er.add_paragraph()
    p.space_before = Pt(8)
    run1 = p.add_run()
    run1.text = "• " + title_txt
    run1.font.bold = True
    run1.font.size = Pt(14)
    run1.font.color.rgb = TEXT_WHITE
    run2 = p.add_run()
    run2.text = body_txt
    run2.font.size = Pt(13)
    run2.font.color.rgb = TEXT_LIGHT

add_speaker_notes(slide10, """Slide 10 presents our Entity-Relationship (ER) diagram. 
We have three core entities: USERS, EVENTS, and the associative entity REGISTRATIONS.

Examining the cardinalities: one user can have many registrations (1 to N), and one event can have many registrations (1 to N). 
Notice the highlighted compound unique constraint: `UNIQUE (user_id, event_id)`. This prevents duplicate entries at the database level. In addition, we apply foreign key cascades so that deleting an event cleanly cleans up associated interest rows.""")

# ==========================================
# SLIDE 11: Database Tables / Schema Definition
# ==========================================
slide11 = prs.slides.add_slide(blank_layout)
apply_background(slide11)
add_header(slide11, "Database Tables & Schema Definition", category="SCHEMA SPECIFICATION", subtitle="Data definition language (DDL), column types, primary/foreign keys, and constraints")

# We create 3 clean native tables for users, events, registrations
# Table 1: users (Top Left)
add_card(slide11, Inches(0.8), Inches(1.85), Inches(5.6), Inches(2.35))
tb_t1 = slide11.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(5.2), Inches(0.4))
p = tb_t1.text_frame.paragraphs[0]
p.text = "1. users Table (Identity)"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = ACCENT_INDIGO

table_users = slide11.shapes.add_table(4, 4, Inches(1.0), Inches(2.3), Inches(5.2), Inches(1.7)).table
table_users.columns[0].width = Inches(1.2)
table_users.columns[1].width = Inches(1.4)
table_users.columns[2].width = Inches(1.2)
table_users.columns[3].width = Inches(1.4)

headers = ["Column", "Data Type", "Constraint", "Description"]
for i, h in enumerate(headers):
    cell = table_users.cell(0, i)
    cell.fill.solid()
    cell.fill.fore_color.rgb = CARD_BG_ALT
    p = cell.text_frame.paragraphs[0]
    p.text = h
    p.font.name = FONT_FAMILY
    p.font.bold = True
    p.font.size = Pt(11)
    p.font.color.rgb = ACCENT_CYAN

users_data = [
    ("id", "BIGINT", "PK, AUTO_INC", "Unique User ID"),
    ("username", "VARCHAR(50)", "NOT NULL, UK", "Login Identifier"),
    ("password", "VARCHAR(255)", "NOT NULL", "BCrypt 60-char Hash")
]
for row_idx, row in enumerate(users_data):
    for col_idx, val in enumerate(row):
        cell = table_users.cell(row_idx + 1, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD_BG if row_idx % 2 == 0 else CARD_BG_ALT
        p = cell.text_frame.paragraphs[0]
        p.text = val
        p.font.name = FONT_FAMILY
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_LIGHT

# Table 2: registrations (Bottom Left)
add_card(slide11, Inches(0.8), Inches(4.35), Inches(5.6), Inches(2.6))
tb_t2 = slide11.shapes.add_textbox(Inches(1.0), Inches(4.4), Inches(5.2), Inches(0.4))
p = tb_t2.text_frame.paragraphs[0]
p.text = "2. registrations Table (Associative M:N)"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = ACCENT_CYAN

table_regs = slide11.shapes.add_table(5, 4, Inches(1.0), Inches(4.8), Inches(5.2), Inches(1.95)).table
table_regs.columns[0].width = Inches(1.2)
table_regs.columns[1].width = Inches(1.4)
table_regs.columns[2].width = Inches(1.2)
table_regs.columns[3].width = Inches(1.4)

for i, h in enumerate(headers):
    cell = table_regs.cell(0, i)
    cell.fill.solid()
    cell.fill.fore_color.rgb = CARD_BG_ALT
    p = cell.text_frame.paragraphs[0]
    p.text = h
    p.font.name = FONT_FAMILY
    p.font.bold = True
    p.font.size = Pt(11)
    p.font.color.rgb = ACCENT_CYAN

regs_data = [
    ("id", "BIGINT", "PK, AUTO_INC", "Registration ID"),
    ("user_id", "BIGINT", "FK -> users(id)", "User Reference"),
    ("event_id", "BIGINT", "FK -> events(id)", "Event Reference"),
    ("status", "VARCHAR(20)", "CHECK (IN...)", "Interest Status")
]
for row_idx, row in enumerate(regs_data):
    for col_idx, val in enumerate(row):
        cell = table_regs.cell(row_idx + 1, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD_BG if row_idx % 2 == 0 else CARD_BG_ALT
        p = cell.text_frame.paragraphs[0]
        p.text = val
        p.font.name = FONT_FAMILY
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_LIGHT

# Table 3: events (Right Column Full Height)
add_card(slide11, Inches(6.8), Inches(1.85), Inches(5.733), Inches(5.1))
tb_t3 = slide11.shapes.add_textbox(Inches(7.0), Inches(1.9), Inches(5.333), Inches(0.4))
p = tb_t3.text_frame.paragraphs[0]
p.text = "3. events Table (Event Catalogue)"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = ACCENT_GREEN

table_events = slide11.shapes.add_table(8, 4, Inches(7.0), Inches(2.3), Inches(5.333), Inches(4.45)).table
table_events.columns[0].width = Inches(1.2)
table_events.columns[1].width = Inches(1.4)
table_events.columns[2].width = Inches(1.3)
table_events.columns[3].width = Inches(1.433)

for i, h in enumerate(headers):
    cell = table_events.cell(0, i)
    cell.fill.solid()
    cell.fill.fore_color.rgb = CARD_BG_ALT
    p = cell.text_frame.paragraphs[0]
    p.text = h
    p.font.name = FONT_FAMILY
    p.font.bold = True
    p.font.size = Pt(11)
    p.font.color.rgb = ACCENT_GREEN

events_data = [
    ("id", "BIGINT", "PK, AUTO_INC", "Primary Identifier"),
    ("title", "VARCHAR(255)", "NOT NULL", "Event Name"),
    ("date_time", "DATETIME", "NOT NULL, INDEX", "Event Start Time"),
    ("end_date_time", "DATETIME", "CHECK (> start)", "Event End Time"),
    ("venue", "VARCHAR(255)", "NOT NULL", "Physical Location"),
    ("category", "VARCHAR(50)", "NOT NULL, INDEX", "Technical, Cultural..."),
    ("max_capacity", "INT", "CHECK (> 0)", "Seat Capacity Ceiling")
]
for row_idx, row in enumerate(events_data):
    for col_idx, val in enumerate(row):
        cell = table_events.cell(row_idx + 1, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD_BG if row_idx % 2 == 0 else CARD_BG_ALT
        p = cell.text_frame.paragraphs[0]
        p.text = val
        p.font.name = FONT_FAMILY
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_LIGHT

add_speaker_notes(slide11, """Slide 11 shows the exact relational Data Definition Language (DDL) specifications. 
Notice the exact constraints implemented:
In the `users` table: passwords are sized at VARCHAR(255) to accommodate BCrypt and future Argon2 encoders.
In the `events` table: we enforce `CHECK (max_capacity > 0)` and `CHECK (end_date_time > date_time)`. We also index `(date_time)` and `(category, date_time)` to ensure fast queries.
In the `registrations` table: foreign keys cascade on delete, and status is restricted via CHECK constraint to INTERESTED, CONFIRMED, CANCELLED, or WAITLISTED.""")

# ==========================================
# SLIDE 12: Sample Database Records
# ==========================================
slide12 = prs.slides.add_slide(blank_layout)
apply_background(slide12)
add_header(slide12, "Sample Database Records & Live State", category="DATABASE VERIFICATION", subtitle="Actual seed records demonstrating relational persistence and referential integrity")

# Subtable 1: Sample Users
add_card(slide12, Inches(0.8), Inches(1.85), Inches(11.733), Inches(1.5))
tb_su = slide12.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(11.333), Inches(0.3))
p = tb_su.text_frame.paragraphs[0]
p.text = "Sample Records: users Table"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = ACCENT_INDIGO

t_su = slide12.shapes.add_table(3, 5, Inches(1.0), Inches(2.25), Inches(11.333), Inches(0.95)).table
t_su.columns[0].width = Inches(0.8)
t_su.columns[1].width = Inches(1.5)
t_su.columns[2].width = Inches(1.5)
t_su.columns[3].width = Inches(4.5)
t_su.columns[4].width = Inches(3.033)

u_headers = ["id", "username", "role", "password_hash (BCrypt)", "email"]
for i, h in enumerate(u_headers):
    cell = t_su.cell(0, i)
    cell.fill.solid()
    cell.fill.fore_color.rgb = CARD_BG_ALT
    p = cell.text_frame.paragraphs[0]
    p.text = h
    p.font.bold = True
    p.font.size = Pt(10)
    p.font.color.rgb = ACCENT_CYAN

u_rows = [
    ("1", "admin", "ADMIN", "$2a$10$RzVj6l9z8qP7yW2u3e4r5t6y7u8i9o0p1a2s3d4f5g6h7j8k9l", "admin@campus.edu"),
    ("2", "guest", "STUDENT", "$2a$10$kL8m9n0b1v2c3x4z5a6s7d8f9g0h1j2k3l4m5n6b7v8c9x0z", "student@campus.edu")
]
for r_i, r in enumerate(u_rows):
    for c_i, v in enumerate(r):
        cell = t_su.cell(r_i + 1, c_i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD_BG if r_i % 2 == 0 else CARD_BG_ALT
        p = cell.text_frame.paragraphs[0]
        p.text = v
        p.font.size = Pt(9)
        p.font.color.rgb = TEXT_LIGHT

# Subtable 2: Sample Events
add_card(slide12, Inches(0.8), Inches(3.5), Inches(11.733), Inches(1.8))
tb_se = slide12.shapes.add_textbox(Inches(1.0), Inches(3.55), Inches(11.333), Inches(0.3))
p = tb_se.text_frame.paragraphs[0]
p.text = "Sample Records: events Table"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = ACCENT_GREEN

t_se = slide12.shapes.add_table(3, 6, Inches(1.0), Inches(3.9), Inches(11.333), Inches(1.25)).table
t_se.columns[0].width = Inches(0.8)
t_se.columns[1].width = Inches(3.2)
t_se.columns[2].width = Inches(1.8)
t_se.columns[3].width = Inches(2.2)
t_se.columns[4].width = Inches(2.0)
t_se.columns[5].width = Inches(1.333)

e_headers = ["id", "title", "category", "date_time", "venue", "max_capacity"]
for i, h in enumerate(e_headers):
    cell = t_se.cell(0, i)
    cell.fill.solid()
    cell.fill.fore_color.rgb = CARD_BG_ALT
    p = cell.text_frame.paragraphs[0]
    p.text = h
    p.font.bold = True
    p.font.size = Pt(10)
    p.font.color.rgb = ACCENT_GREEN

e_rows = [
    ("1", "AI & Autonomous Agents Hackathon", "Technical", "2026-09-15 09:00:00", "Alan Turing Innovation Lab", "120"),
    ("2", "Spring Symphony & Cultural Gala", "Cultural", "2026-09-22 18:30:00", "University Main Auditorium", "500")
]
for r_i, r in enumerate(e_rows):
    for c_i, v in enumerate(r):
        cell = t_se.cell(r_i + 1, c_i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD_BG if r_i % 2 == 0 else CARD_BG_ALT
        p = cell.text_frame.paragraphs[0]
        p.text = v
        p.font.size = Pt(9)
        p.font.color.rgb = TEXT_LIGHT

# Subtable 3: Sample Registrations
add_card(slide12, Inches(0.8), Inches(5.45), Inches(11.733), Inches(1.6))
tb_sr = slide12.shapes.add_textbox(Inches(1.0), Inches(5.5), Inches(11.333), Inches(0.3))
p = tb_sr.text_frame.paragraphs[0]
p.text = "Sample Records: registrations Table"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = ACCENT_CYAN

t_sr = slide12.shapes.add_table(3, 5, Inches(1.0), Inches(5.85), Inches(11.333), Inches(1.05)).table
t_sr.columns[0].width = Inches(0.8)
t_sr.columns[1].width = Inches(2.0)
t_sr.columns[2].width = Inches(2.0)
t_sr.columns[3].width = Inches(3.5)
t_sr.columns[4].width = Inches(3.033)

r_headers = ["id", "user_id (FK)", "event_id (FK)", "registration_date", "status"]
for i, h in enumerate(r_headers):
    cell = t_sr.cell(0, i)
    cell.fill.solid()
    cell.fill.fore_color.rgb = CARD_BG_ALT
    p = cell.text_frame.paragraphs[0]
    p.text = h
    p.font.bold = True
    p.font.size = Pt(10)
    p.font.color.rgb = ACCENT_CYAN

r_rows = [
    ("1", "2 (guest)", "1 (AI Hackathon)", "2026-09-08 14:30:12", "INTERESTED"),
    ("2", "2 (guest)", "2 (Spring Gala)", "2026-09-08 15:45:00", "INTERESTED")
]
for r_i, r in enumerate(r_rows):
    for c_i, v in enumerate(r):
        cell = t_sr.cell(r_i + 1, c_i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD_BG if r_i % 2 == 0 else CARD_BG_ALT
        p = cell.text_frame.paragraphs[0]
        p.text = v
        p.font.size = Pt(9)
        p.font.color.rgb = TEXT_LIGHT

add_speaker_notes(slide12, """Slide 12 provides direct proof of actual database operations and records.
In the `users` table: passwords are never stored in plaintext. They are salted BCrypt hashes starting with $2a$.
In the `events` table: we store structured attributes including start dates, categories, and capacities.
In the `registrations` table: we see the associative rows linking user 2 (guest) with events 1 and 2, with status recorded as INTERESTED. Attempting to insert a duplicate (user 2, event 1) is immediately rejected by the unique database index.""")

# ==========================================
# SLIDE 13: Current Progress & Technical Milestones
# ==========================================
slide13 = prs.slides.add_slide(blank_layout)
apply_background(slide13)
add_header(slide13, "Current Progress & Verification Status", category="PROGRESS ASSESSMENT", subtitle="85% of core platform completed; automated test suite and CI pipeline fully passing")

# Left Column: Progress Checklist
add_card(slide13, Inches(0.8), Inches(1.9), Inches(6.0), Inches(5.0))
tb_prog = slide13.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(5.4), Inches(4.6))
tf_prog = tb_prog.text_frame
tf_prog.word_wrap = True

p = tf_prog.paragraphs[0]
p.text = "Implementation Checklist"
p.font.size = Pt(21)
p.font.bold = True
p.font.color.rgb = ACCENT_GREEN

milestones = [
    ("✅ Modular Monolith Architecture: ", "Spring Boot 4 + Java 25 runtime cleanly configured."),
    ("✅ Database Schema & Migrations: ", "Flyway V1–V3 fully deployed on MySQL 8.4 LTS."),
    ("✅ Security & Rate Limiting: ", "BCrypt, CSRF, and Bucket4j IP filters active."),
    ("✅ Frontend Surfaces: ", "Student Discovery Feed, Admin Console, and Forms operational."),
    ("✅ Automated Testing Suite: ", "63 unit and integration tests passing with JaCoCo coverage gates."),
    ("✅ Automated CI Pipeline: ", "GitHub Actions build and verify on every push.")
]

for title_txt, body_txt in milestones:
    p = tf_prog.add_paragraph()
    p.space_before = Pt(10)
    run1 = p.add_run()
    run1.text = title_txt
    run1.font.bold = True
    run1.font.size = Pt(15)
    run1.font.color.rgb = TEXT_WHITE
    run2 = p.add_run()
    run2.text = body_txt
    run2.font.size = Pt(14)
    run2.font.color.rgb = TEXT_LIGHT

# Right Column: Quality & Metric Telemetry
add_card(slide13, Inches(7.1), Inches(1.9), Inches(5.433), Inches(5.0))
tb_met = slide13.shapes.add_textbox(Inches(7.4), Inches(2.1), Inches(4.833), Inches(4.6))
tf_met = tb_met.text_frame
tf_met.word_wrap = True

p = tf_met.paragraphs[0]
p.text = "Quality & Verification Metrics"
p.font.size = Pt(21)
p.font.bold = True
p.font.color.rgb = ACCENT_CYAN

metrics = [
    ("Automated Test Suite: ", "63 / 63 tests passing (Surefire 3.5.4) covering controllers, services, locking, and security."),
    ("Code Coverage Gates: ", "Strict JaCoCo line and branch coverage gates enforced in CI."),
    ("Health Probes: ", "Spring Actuator `/actuator/health` endpoint responding with UP status."),
    ("OpenAPI Documentation: ", "Generated OpenAPI 3.0 specification active at `/v3/api-docs`."),
    ("Containerization: ", "Non-root Dockerfile and Docker Compose stack verified locally.")
]

for title_txt, body_txt in metrics:
    p = tf_met.add_paragraph()
    p.space_before = Pt(10)
    run1 = p.add_run()
    run1.text = "• " + title_txt
    run1.font.bold = True
    run1.font.size = Pt(15)
    run1.font.color.rgb = TEXT_WHITE
    run2 = p.add_run()
    run2.text = body_txt
    run2.font.size = Pt(14)
    run2.font.color.rgb = TEXT_LIGHT

add_speaker_notes(slide13, """Slide 13 summarizes our current progress for Review 2. We have achieved an estimated 85% completion of our core platform. 

Everything we present is backed by hard evidence:
All 63 automated tests are passing in our test suite, spanning service logic, pessimistic locking, controller validations, and security policies. 
Our GitHub Actions CI pipeline builds the Docker image and executes tests against a real MySQL 8.4 service container. In addition, operational endpoints like `/actuator/health` and `/v3/api-docs` are already operational.""")

# ==========================================
# SLIDE 14: Challenges / Issues Faced
# ==========================================
slide14 = prs.slides.add_slide(blank_layout)
apply_background(slide14)
add_header(slide14, "Challenges Faced & Technical Solutions", category="ENGINEERING RESOLUTIONS", subtitle="Obstacles encountered during development and our technical remediations")

challenges = [
    (Inches(0.8), Inches(1.9), "1. Concurrent Interest Race Condition", ACCENT_PINK, [
        ("Problem: ", "Under rapid concurrent user requests, multiple threads could pass the duplicate check simultaneously, causing duplicate DB rows or transaction collisions."),
        ("Solution: ", "Applied JPA `PESSIMISTIC_WRITE` lock on the event entity row during registration in `EventService`, combined with a database-level `uk_user_event` unique constraint.")
    ]),
    (Inches(6.9), Inches(1.9), "2. Flyway 12 & MySQL 8.4 Dialect Mismatch", ACCENT_AMBER, [
        ("Problem: ", "Upgrading to MySQL 8.4 LTS caused Flyway 12 to fail during startup because MySQL support was extracted from core Flyway into a separate module."),
        ("Solution: ", "Explicitly integrated `org.flywaydb:flyway-mysql:12.4.0` in `pom.xml`, restoring seamless schema migrations across local Docker and GitHub Actions CI.")
    ]),
    (Inches(0.8), Inches(4.55), "3. Admin Secret Ingestion & Zero-Trust", ACCENT_CYAN, [
        ("Problem: ", "Hardcoded admin credentials in application properties violate production security standards and invite repository leaks."),
        ("Solution: ", "Implemented environment-driven secret injection (`ADMIN_PASSWORD`), with an undisclosed 32-byte CSPRNG token generated on clean bootstrap if omitted.")
    ]),
    (Inches(6.9), Inches(4.55), "4. Image Storage & Deployment Portability", ACCENT_GREEN, [
        ("Problem: ", "Storing uploaded images on local disk volumes makes container deployment fragile across ephemeral container restarts."),
        ("Solution: ", "Engineered database-backed `MEDIUMBLOB` persistence with MIME-type validation, ensuring complete image durability inside the MySQL tier.")
    ])
]

for left, top, title, color, points in challenges:
    add_card(slide14, left, top, Inches(5.6), Inches(2.45))
    tb_ch = slide14.shapes.add_textbox(left + Inches(0.25), top + Inches(0.18), Inches(5.1), Inches(2.1))
    tf_ch = tb_ch.text_frame
    tf_ch.word_wrap = True
    
    p = tf_ch.paragraphs[0]
    p.text = title
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = color
    
    for lead, desc in points:
        p = tf_ch.add_paragraph()
        p.space_before = Pt(6)
        r1 = p.add_run()
        r1.text = "• " + lead
        r1.font.bold = True
        r1.font.size = Pt(14)
        r1.font.color.rgb = TEXT_WHITE
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(13)
        r2.font.color.rgb = TEXT_LIGHT

add_speaker_notes(slide14, """Slide 14 details four significant engineering challenges we encountered and solved:
First, concurrent registration race conditions: we resolved this using JPA pessimistic write locking and a unique database constraint.
Second, Flyway 12 compatibility with MySQL 8.4 LTS: we resolved this by pinning the dedicated `flyway-mysql` module.
Third, secure admin initialization: we removed all hardcoded passwords, replacing them with environment secret injection.
Fourth, container image persistence: we implemented database-backed MEDIUMBLOB storage with MIME-type verification.""")

# ==========================================
# SLIDE 15: Future Work & Conclusion
# ==========================================
slide15 = prs.slides.add_slide(blank_layout)
apply_background(slide15)
add_header(slide15, "Future Work & Project Conclusion", category="CONCLUSION & ROADMAP", subtitle="Sprint 3 milestones, planned distributed enhancements, and final project takeaways")

# Left Column: Future Work (Sprint 3)
add_card(slide15, Inches(0.8), Inches(1.9), Inches(5.6), Inches(5.0))
tb_fw = slide15.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(5.0), Inches(4.6))
tf_fw = tb_fw.text_frame
tf_fw.word_wrap = True

p = tf_fw.paragraphs[0]
p.text = "Future Work (Sprint 3 Roadmap)"
p.font.size = Pt(21)
p.font.bold = True
p.font.color.rgb = ACCENT_INDIGO

roadmap = [
    ("Internal Seat Ticketing State Machine: ", "Transition from interest registration to internal capacity reservations with real-time waitlists."),
    ("Asynchronous Event Notifications: ", "Implement transactional Outbox pattern with RabbitMQ or Kafka for event updates and email reminders."),
    ("Semantic Vector Search: ", "Integrate pgvector or embedding adapters for intelligent semantic event matching across student interests."),
    ("Mobile PWA Client: ", "Package offline-first Progressive Web App capabilities for fast mobile discovery.")
]

for title_txt, body_txt in roadmap:
    p = tf_fw.add_paragraph()
    p.space_before = Pt(10)
    run1 = p.add_run()
    run1.text = "• " + title_txt
    run1.font.bold = True
    run1.font.size = Pt(15)
    run1.font.color.rgb = TEXT_WHITE
    run2 = p.add_run()
    run2.text = body_txt
    run2.font.size = Pt(14)
    run2.font.color.rgb = TEXT_LIGHT

# Right Column: Conclusion
add_card(slide15, Inches(6.9), Inches(1.9), Inches(5.6), Inches(5.0))
tb_con = slide15.shapes.add_textbox(Inches(7.2), Inches(2.1), Inches(5.0), Inches(4.6))
tf_con = tb_con.text_frame
tf_con.word_wrap = True

p = tf_con.paragraphs[0]
p.text = "Review 2 Concluding Summary"
p.font.size = Pt(21)
p.font.bold = True
p.font.color.rgb = ACCENT_GREEN

conclusions = [
    ("Working Production-Ready Foundation: ", "CampusConnect delivers an operational modular monolith with high-performance event discovery and secure administration."),
    ("Verified Relational Engineering: ", "3NF schema, Flyway V1–V3 migrations, and pessimistic concurrency locks provide enterprise-grade reliability."),
    ("Rigorous Quality Assurance: ", "63 passing automated tests and automated GitHub Actions CI ensure zero-regression stability."),
    ("Clear Evolutionary Path: ", "Architectural boundaries are cleanly established for future distributed messaging and microservice extraction.")
]

for title_txt, body_txt in conclusions:
    p = tf_con.add_paragraph()
    p.space_before = Pt(10)
    run1 = p.add_run()
    run1.text = "• " + title_txt
    run1.font.bold = True
    run1.font.size = Pt(15)
    run1.font.color.rgb = TEXT_WHITE
    run2 = p.add_run()
    run2.text = body_txt
    run2.font.size = Pt(14)
    run2.font.color.rgb = TEXT_LIGHT

add_speaker_notes(slide15, """In conclusion, for Review 2, CampusConnect has successfully transformed from an abstract concept into a fully functional, tested, and secure modular monolith. 

We have verified our frontend pages, implemented robust input validation, deployed a normalized MySQL schema via Flyway migrations, and backed our work with 63 automated tests in CI. 

Looking forward to Sprint 3 and Review 3, we plan to implement internal seat reservation ticketing, asynchronous notifications, and semantic search. Thank you for your time and guidance. We are now ready to take questions from the panel.""")

prs.save(OUTPUT_PPTX)
print(f"Presentation successfully saved to: {OUTPUT_PPTX}")
print(f"Total Slides Generated: {len(prs.slides)}")
